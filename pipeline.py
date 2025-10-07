# pipeline.py
import os
import cv2
import pytesseract
import pdfplumber
from pathlib import Path
from PIL import Image
from ultralytics import YOLO
from presidio_analyzer import AnalyzerEngine, Pattern, PatternRecognizer
from presidio_anonymizer import AnonymizerEngine
from dotenv import load_dotenv
import google.generativeai as genai
import openpyxl
from openpyxl import load_workbook
from pptx import Presentation
from datetime import datetime

# --- Configuration ---
UPLOAD_DIR = Path("uploads")
REDACTED_DIR = Path("redacted_files")
MODEL_DIR = Path("models")

# Ensure directories exist
UPLOAD_DIR.mkdir(exist_ok=True)
REDACTED_DIR.mkdir(exist_ok=True)

# Load environment variables
load_dotenv()

# Configure Gemini API
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print("Warning: GEMINI_API_KEY not found in environment variables")
else:
    genai.configure(api_key=GEMINI_API_KEY)

# --- Global Variables for Models (initialized on first use) ---
_yolo_model = None
_gemini_model = None
_analyzer = None
_anonymizer = None

# --- Model Loading Functions (Lazy Loading) ---
def get_yolo_model():
    """Load YOLO model (lazy loading)"""
    global _yolo_model
    if (_yolo_model is None):
        YOLO_MODEL_PATH = "best.pt"
        if not Path(YOLO_MODEL_PATH).exists():
            print(f"Warning: YOLO model not found at {YOLO_MODEL_PATH}")
            return None
        _yolo_model = YOLO(YOLO_MODEL_PATH)
    return _yolo_model

def get_gemini_model():
    """Get Gemini model (lightweight)"""
    global _gemini_model
    if _gemini_model is None:
        try:
            _gemini_model = genai.GenerativeModel('gemini-2.5-flash')
            print("Gemini model initialized successfully")
        except Exception as e:
            print(f"Warning: Could not initialize Gemini model: {e}")
            return None
    return _gemini_model

def get_analyzer():
    """Get Presidio analyzer with comprehensive PII and technical data recognizers"""
    global _analyzer
    if (_analyzer is None):
        # AWS and Cloud Infrastructure Patterns
        aws_patterns = [
            #Pattern(name="AWS_ARN", regex=r"arn:aws:[a-z0-9\-]+:[a-z0-9\-]*:[0-9]{12}:[a-zA-Z0-9\-\/\:\_]+", score=0.65),
            #Pattern(name="AWS_ROLE_ID", regex=r"A[IKR][A-Z0-9]{18,}", score=0.60),
            Pattern(name="AWS_ACCOUNT", regex=r"\b\d{12}\b", score=0.85),
            Pattern(name="IP_CIDR", regex=r"\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}/\d{1,2}\b", score=0.60),
            Pattern(name="IP_ADDRESS", regex=r"\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b", score=0.98),
            Pattern(name="MAC_ADDRESS", regex=r"\b([0-9A-Fa-f]{2}[:-]){5}([0-9A-Fa-f]{2})\b", score=0.95),
        ]
        
        # Network and Security Patterns
        network_patterns = [
            Pattern(name="PORT_RANGE", regex=r"\b(tcp|udp|icmp):\d{1,5}(-\d{1,5})?\b", score=0.95),
            Pattern(name="PORT_SINGLE", regex=r"\b(port|Port)\s*:?\s*\d{1,5}\b", score=0.98),
            Pattern(name="PROTOCOL", regex=r"\b(tcp|udp|icmp|ssh|rdp|sftp):\d+\b", score=0.85),
        ]
        
        # File paths and sensitive locations
        path_patterns = [
            Pattern(name="FILE_PATH", regex=r"file://[a-zA-Z0-9\-\_\./]+", score=0.74),
            Pattern(name="UNIX_PATH", regex=r"/[a-zA-Z0-9\-\_/]+\.[a-zA-Z]{2,4}", score=0.7),
        ]
        # Name Patterns - Enhanced for better detection
        name_patterns = [
            # Full names (First Last, First Middle Last)
            Pattern(name="FULL_NAME", regex=r"\b[A-Z][a-z]{2,}\s+[A-Z][a-z]{2,}(?:\s+[A-Z][a-z]{2,})?\b", score=0.85),
            # Names with titles (Mr., Mrs., Dr., etc.)
            Pattern(name="TITLED_NAME", regex=r"\b(?:Mr\.?|Mrs\.?|Ms\.?|Dr\.?|Prof\.?|Sir|Madam)\s+[A-Z][a-z]{2,}(?:\s+[A-Z][a-z]{2,})?\b", score=0.90),
            # Names in "Last, First" format
            Pattern(name="LASTNAME_FIRST", regex=r"\b[A-Z][a-z]{2,},\s+[A-Z][a-z]{2,}(?:\s+[A-Z]\.?)?\b", score=0.88),
            # Employee/Person IDs with names
            Pattern(name="EMPLOYEE_NAME", regex=r"(?:Employee|Staff|User|Person):\s*[A-Z][a-z]{2,}\s+[A-Z][a-z]{2,}", score=0.90),
            # Names in brackets or quotes
            Pattern(name="QUOTED_NAME", regex=r"[\"']([A-Z][a-z]{2,}\s+[A-Z][a-z]{2,})[\"']", score=0.85),
            # Names with common suffixes
            Pattern(name="NAME_WITH_SUFFIX", regex=r"\b[A-Z][a-z]{2,}\s+[A-Z][a-z]{2,}\s+(?:Jr\.?|Sr\.?|III?|IV)\b", score=0.90),
        ]

         # Date Patterns - Comprehensive date formats
        date_patterns = [
            # MM/DD/YYYY, MM-DD-YYYY, MM.DD.YYYY
            Pattern(name="US_DATE_FORMAT", regex=r"\b(?:0?[1-9]|1[0-2])[\/\-\.]\d{1,2}[\/\-\.]\d{4}\b", score=0.85),
            # DD/MM/YYYY, DD-MM-YYYY, DD.MM.YYYY
            Pattern(name="EU_DATE_FORMAT", regex=r"\b(?:0?[1-9]|[12]\d|3[01])[\/\-\.]\d{1,2}[\/\-\.]\d{4}\b", score=0.80),
            # YYYY-MM-DD, YYYY/MM/DD (ISO format)
            Pattern(name="ISO_DATE_FORMAT", regex=r"\b\d{4}[\/\-\.]\d{1,2}[\/\-\.]\d{1,2}\b", score=0.90),
            # Month DD, YYYY
            Pattern(name="WRITTEN_DATE", regex=r"\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4}\b", score=0.85),
            # DD Month YYYY
            Pattern(name="WRITTEN_DATE_EU", regex=r"\b\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b", score=0.85),
            # Short month formats (Jan, Feb, etc.)
            Pattern(name="SHORT_MONTH_DATE", regex=r"\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\.?\s+\d{1,2},?\s+\d{4}\b", score=0.80),
            # Date with time (YYYY-MM-DD HH:MM:SS)
            Pattern(name="DATETIME_ISO", regex=r"\b\d{4}-\d{2}-\d{2}\s+\d{2}:\d{2}:\d{2}\b", score=0.90),
            # Timestamp formats
            Pattern(name="TIMESTAMP", regex=r"\b\d{2}\/\d{2}\/\d{4}\s+\d{1,2}:\d{2}(?::\d{2})?\s*(?:AM|PM)?\b", score=0.85),
            # Birth dates (common formats)
            Pattern(name="BIRTH_DATE", regex=r"(?:DOB|Date of Birth|Born):\s*\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{4}", score=0.95),
            # Expiry dates
            Pattern(name="EXPIRY_DATE", regex=r"(?:Exp|Expiry|Expires?):\s*\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4}", score=0.90),
            # Relative dates
            Pattern(name="RELATIVE_DATE", regex=r"\b(?:yesterday|today|tomorrow|\d+\s+(?:days?|weeks?|months?|years?)\s+(?:ago|from now))\b", score=0.75),
        ]

        # Additional Person-related Patterns
        person_patterns = [
            # Social media handles
            Pattern(name="SOCIAL_HANDLE", regex=r"@[A-Za-z0-9_]{3,}", score=0.70),
            # Initials (A.B., A.B.C.)
            Pattern(name="INITIALS", regex=r"\b[A-Z]\.[A-Z]\.(?:[A-Z]\.)?\b", score=0.75),
            # Signatures or "Signed by"
            Pattern(name="SIGNATURE", regex=r"(?:Signed by|Signature|Regards|Best regards|Sincerely),?\s*[A-Z][a-z]{2,}(?:\s+[A-Z][a-z]{2,})?", score=0.80),
        ]

        # Create recognizers
        aws_recognizer = PatternRecognizer(
            supported_entity="AWS_IDENTIFIER",
            patterns=aws_patterns,
            name="AWS Recognizer"
        )
        
        network_recognizer = PatternRecognizer(
            supported_entity="NETWORK_INFO",
            patterns=network_patterns,
            name="Network Recognizer"
        )
        
        path_recognizer = PatternRecognizer(
            supported_entity="FILE_PATH",
            patterns=path_patterns,
            name="Path Recognizer"
        )
         # New recognizers for names and dates
        name_recognizer = PatternRecognizer(
            supported_entity="PERSON_NAME",
            patterns=name_patterns,
            name="Enhanced Name Recognizer"
        )
        
        date_recognizer = PatternRecognizer(
            supported_entity="DATE_TIME",
            patterns=date_patterns,
            name="Enhanced Date Recognizer"
        )
        
        person_recognizer = PatternRecognizer(
            supported_entity="PERSON_INFO",
            patterns=person_patterns,
            name="Person Info Recognizer"
        )
        
        # Initialize analyzer with all recognizers
        _analyzer = AnalyzerEngine()
        _analyzer.registry.add_recognizer(aws_recognizer)
        _analyzer.registry.add_recognizer(network_recognizer)
        _analyzer.registry.add_recognizer(path_recognizer)
        _analyzer.registry.add_recognizer(name_recognizer)
        _analyzer.registry.add_recognizer(date_recognizer)
        _analyzer.registry.add_recognizer(person_recognizer)
    
    return _analyzer

def get_anonymizer():
    """Get Presidio anonymizer"""
    global _anonymizer
    if (_anonymizer is None):
        _anonymizer = AnonymizerEngine()
    return _anonymizer

# --- Helper Functions ---
def extract_file_info_from_path(file_path):
    """Extract file name and type from path for table formatting"""
    if file_path is None:
        return {'name': 'Document', 'type': '.png', 'full_name': 'Document.png'}
    
    path = Path(file_path)
    # Extract base name without _redacted suffix
    name = path.stem.replace('_redacted_llava', '').replace('_redacted_text', '').replace('_redacted_scenic', '')
    return {
        'name': name,
        'type': path.suffix,
        'full_name': path.name
    }

# Update the describe_image_with_gemini function
def describe_image_with_gemini(file_path: Path):
    """Use Gemini to describe what's in the image"""
    try:
        gemini_model = get_gemini_model()
        if gemini_model is None:
            return "Image contains visual content. No text detected for PII analysis."
        
        # Load and prepare image
        img = Image.open(file_path)
        
        # Resize image if too large for API
        max_size = 1024
        if img.width > max_size or img.height > max_size:
            img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
        
        prompt = """Analyze this image and provide a brief description. Focus on:
1. What type of content, document, or scene this is
2. Main objects, elements, or information visible
3. The apparent purpose or context

Provide 2-3 clear sentences. Be specific about what you see."""
        
        response = gemini_model.generate_content([prompt, img])
        description = response.text.strip()
        
        if not description or len(description) < 10:
            return "Image processed successfully. Visual content detected but no detailed description available."
        
        return description
        
    except Exception as e:
        print(f"Error describing image with Gemini: {e}")
        return "Image contains visual content. Unable to generate detailed description due to processing limitations."

# Add this new function after the describe_image_with_gemini function

def detect_and_remove_logos(file_path: Path, img):
    """Use Gemini to detect logos and remove them from the image"""
    try:
        gemini_model = get_gemini_model()
        if gemini_model is None:
            print("Gemini model not available for logo detection")
            return img, []
        
        # Convert OpenCV image to PIL for Gemini
        img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(img_rgb)
        
        # Resize if too large
        max_size = 1024
        if pil_img.width > max_size or pil_img.height > max_size:
            pil_img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
        
        prompt = """Analyze this image and identify any logos, brand marks, or company identifiers. 

For each logo you find, provide:
1. A brief description of the logo
2. The approximate location (top-left, top-right, center, bottom-left, bottom-right, etc.)
3. Whether it's a company logo, brand mark, or watermark

Format your response as:
LOGO_FOUND: [description] | LOCATION: [location] | TYPE: [company/brand/watermark]

If no logos are found, respond with: NO_LOGOS_DETECTED

Examples:
LOGO_FOUND: Microsoft logo with four colored squares | LOCATION: top-left | TYPE: company
LOGO_FOUND: Nike swoosh symbol | LOCATION: bottom-right | TYPE: brand
LOGO_FOUND: Watermark text "CONFIDENTIAL" | LOCATION: center | TYPE: watermark"""
        
        response = gemini_model.generate_content([prompt, pil_img])
        detection_result = response.text.strip()
        
        print(f"Logo detection result: {detection_result}")
        
        if "NO_LOGOS_DETECTED" in detection_result.upper():
            return img, []
        
        # Parse logo detections
        detected_logos = []
        lines = detection_result.split('\n')
        
        for line in lines:
            if "LOGO_FOUND:" in line.upper():
                try:
                    parts = line.split('|')
                    description = parts[0].split('LOGO_FOUND:')[1].strip()
                    location = parts[1].split('LOCATION:')[1].strip() if len(parts) > 1 else "unknown"
                    logo_type = parts[2].split('TYPE:')[1].strip() if len(parts) > 2 else "unknown"
                    
                    detected_logos.append({
                        'description': description,
                        'location': location,
                        'type': logo_type
                    })
                except:
                    continue
        
        if not detected_logos:
            return img, []
        
        # Create logo-removed image
        logo_removed_img = remove_logos_from_image(img, detected_logos, pil_img)
        
        return logo_removed_img, detected_logos
        
    except Exception as e:
        print(f"Error in logo detection: {e}")
        return img, []

def remove_logos_from_image(img, detected_logos, pil_img):
    """Remove detected logos from the image using region masking"""
    try:
        h, w = img.shape[:2]
        logo_removed_img = img.copy()
        
        for logo in detected_logos:
            location = logo['location'].lower()
            
            # Define regions based on location descriptions
            if 'top-left' in location or 'top left' in location:
                region = (0, 0, w//3, h//3)
            elif 'top-right' in location or 'top right' in location:
                region = (2*w//3, 0, w, h//3)
            elif 'bottom-left' in location or 'bottom left' in location:
                region = (0, 2*h//3, w//3, h)
            elif 'bottom-right' in location or 'bottom right' in location:
                region = (2*w//3, 2*h//3, w, h)
            elif 'top' in location:
                region = (w//4, 0, 3*w//4, h//4)
            elif 'bottom' in location:
                region = (w//4, 3*h//4, 3*w//4, h)
            elif 'left' in location:
                region = (0, h//4, w//4, 3*h//4)
            elif 'right' in location:
                region = (3*w//4, h//4, w, 3*h//4)
            elif 'center' in location or 'middle' in location:
                region = (w//3, h//3, 2*w//3, 2*h//3)
            else:
                # Default to center region if location unclear
                region = (w//3, h//3, 2*w//3, 2*h//3)
            
            x1, y1, x2, y2 = region
            
            # Apply region-specific logo removal
            logo_removed_img = apply_logo_removal(logo_removed_img, x1, y1, x2, y2, logo)
        
        return logo_removed_img
        
    except Exception as e:
        print(f"Error removing logos: {e}")
        return img

def apply_logo_removal(img, x1, y1, x2, y2, logo_info):
    """Apply logo removal techniques to a specific region"""
    try:
        # Extract the region
        region = img[y1:y2, x1:x2]
        
        if region.size == 0:
            return img
        
        logo_type = logo_info['type'].lower()
        
        if 'watermark' in logo_type or 'text' in logo_info['description'].lower():
            # For text watermarks, use morphological operations
            processed_region = remove_text_watermark(region)
        else:
            # For logos and brand marks, use inpainting
            processed_region = remove_logo_inpainting(region)
        
        # Replace the region in the original image
        img[y1:y2, x1:x2] = processed_region
        
        return img
        
    except Exception as e:
        print(f"Error in logo removal application: {e}")
        return img

def remove_text_watermark(region):
    """Remove text watermarks using morphological operations"""
    try:
        # Convert to grayscale
        gray = cv2.cvtColor(region, cv2.COLOR_BGR2GRAY)
        
        # Create morphological kernel
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (3, 3))
        
        # Apply morphological opening to remove small text elements
        opened = cv2.morphologyEx(gray, cv2.MORPH_OPEN, kernel)
        
        # Create mask for text areas
        text_mask = cv2.absdiff(gray, opened)
        text_mask = cv2.threshold(text_mask, 10, 255, cv2.THRESH_BINARY)[1]
        
        # Dilate to cover text completely
        text_mask = cv2.dilate(text_mask, kernel, iterations=2)
        
        # Inpaint the text areas
        result = cv2.inpaint(region, text_mask, 3, cv2.INPAINT_TELEA)
        
        return result
        
    except Exception as e:
        print(f"Error in text watermark removal: {e}")
        return region

def remove_logo_inpainting(region):
    """Remove logos using inpainting techniques"""
    try:
        # Convert to grayscale for edge detection
        gray = cv2.cvtColor(region, cv2.COLOR_BGR2GRAY)
        
        # Apply Gaussian blur to reduce noise
        blurred = cv2.GaussianBlur(gray, (5, 5), 0)
        
        # Detect edges (potential logo boundaries)
        edges = cv2.Canny(blurred, 50, 150)
        
        # Dilate edges to create mask
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (5, 5))
        mask = cv2.dilate(edges, kernel, iterations=2)
        
        # Apply median filter to reduce small edge artifacts
        mask = cv2.medianBlur(mask, 5)
        
        # Inpaint the logo areas
        result = cv2.inpaint(region, mask, 5, cv2.INPAINT_NS)
        
        return result
        
    except Exception as e:
        print(f"Error in logo inpainting: {e}")
        return region

# --- Main Processing Functions ---
def process_image(file_path: Path):
    """Main function to process images - simplified to just use OCR"""
    print("Processing image with OCR...")
    return process_text_image(file_path, get_analyzer())

# Update the process_text_image function
def process_text_image(file_path: Path, analyzer_engine: AnalyzerEngine):
    """Process text-heavy images with OCR + PII redaction + Logo removal"""
    print(f"OCR-based processing on {file_path.name}...")
    
    img = cv2.imread(str(file_path))
    if img is None:
        return None, "Error: Could not read image"
    
    MAX_DIMENSION = 1800
    h, w, _ = img.shape
    if h > MAX_DIMENSION or w > MAX_DIMENSION:
        scale = MAX_DIMENSION / max(h, w)
        new_w, new_h = int(w * scale), int(h * scale)
        img = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_AREA)
    
    # Step 1: Detect and remove logos
    print("Detecting and removing logos...")
    logo_removed_img, detected_logos = detect_and_remove_logos(file_path, img)
    
    # Log logo detection results
    if detected_logos:
        print(f"Detected {len(detected_logos)} logo(s):")
        for i, logo in enumerate(detected_logos, 1):
            print(f"  {i}. {logo['description']} at {logo['location']} (type: {logo['type']})")
    else:
        print("No logos detected")
    
    # Step 2: Continue with OCR on logo-removed image
    gray = cv2.cvtColor(logo_removed_img, cv2.COLOR_BGR2GRAY)
    gray = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                  cv2.THRESH_BINARY, 11, 2)
    
    ocr_config = r'--oem 3 --psm 6'
    ocr_data = pytesseract.image_to_data(gray, output_type=pytesseract.Output.DICT, config=ocr_config)
    
    full_text_for_analysis = " ".join(word for word in ocr_data['text'] if word.strip())
    
    # If no significant text detected, describe the image content
    if not full_text_for_analysis or len(full_text_for_analysis.strip()) < 20:
        print("No significant text detected, generating image description...")
        image_description = describe_image_with_gemini(file_path)
        
        # Create a redacted image (logo-removed version)
        redacted_filename = f"{file_path.stem}_redacted_text.png"
        redacted_path = REDACTED_DIR / redacted_filename
        cv2.imwrite(str(redacted_path), logo_removed_img)
        
        # Format the description as analysis text, including logo information
        logo_info = ""
        if detected_logos:
            logo_descriptions = [f"{logo['description']} ({logo['location']})" for logo in detected_logos]
            logo_info = f" Logos detected and removed: {', '.join(logo_descriptions)}."
        
        analysis_text = f"Visual Content Analysis: {image_description}. This image does not contain significant readable text requiring PII redaction.{logo_info}"
        
        return redacted_path, analysis_text
    
    # Step 3: Continue with normal PII processing if text is found
    analyzer_results = analyzer_engine.analyze(
        text=full_text_for_analysis, 
        language='en',
        entities=[
            "PERSON", "EMAIL_ADDRESS", "PHONE_NUMBER", "LOCATION",
            "CREDIT_CARD", "US_SSN", "US_DRIVER_LICENSE", "US_PASSPORT",
            "DATE_TIME", "NRP", "MEDICAL_LICENSE", "URL",
            "AWS_IDENTIFIER", "NETWORK_INFO", "FILE_PATH", "TIMESTAMP",
            "IDENTIFIER_NAME", "DURATION","PERSON_NAME","PERSON_INFO"
        ]
    )
    
    pii_char_indices = set()
    for res in analyzer_results:
        for i in range(res.start, res.end):
            pii_char_indices.add(i)
    
    # Apply PII redaction on the logo-removed image
    redacted_img = logo_removed_img.copy()
    current_char_index = 0
    
    for i, word in enumerate(ocr_data['text']):
        if word.strip():
            word_len = len(word)
            conf = int(ocr_data['conf'][i]) if ocr_data['conf'][i] != '-1' else 0
            
            is_pii = any((current_char_index + char_pos) in pii_char_indices 
                        for char_pos in range(word_len))
            
            if is_pii or (conf > 0 and conf < 60 and len(word) > 3):
                (x, y, w, h) = (ocr_data['left'][i], ocr_data['top'][i], 
                               ocr_data['width'][i], ocr_data['height'][i])
                
                padding = 5
                x = max(0, x - padding)
                y = max(0, y - padding)
                w = min(redacted_img.shape[1] - x, w + 2*padding)
                h = min(redacted_img.shape[0] - y, h + 2*padding)
                
                roi = redacted_img[y:y+h, x:x+w]
                if roi.size > 0:
                    cv2.rectangle(redacted_img, (x, y), (x+w, y+h), (0, 0, 0), -1)
            
            current_char_index += word_len + 1
    
    redacted_filename = f"{file_path.stem}_redacted_text.png"
    redacted_path = REDACTED_DIR / redacted_filename
    cv2.imwrite(str(redacted_path), redacted_img)
    
    # Add logo information to analysis text
    logo_info = ""
    if detected_logos:
        logo_descriptions = [f"{logo['description']} ({logo['location']})" for logo in detected_logos]
        logo_info = f"\n\nLogos detected and removed: {', '.join(logo_descriptions)}"
    
    enhanced_text = full_text_for_analysis + logo_info
    
    print(f"Redacted image saved to: {redacted_path}")
    print(f"Extracted text length: {len(full_text_for_analysis)}")
    if detected_logos:
        print(f"Logos removed: {len(detected_logos)}")
    
    return redacted_path, enhanced_text

def process_pdf(file_path: Path):
    """Process PDF files with proper text extraction and analysis"""
    print(f"Processing PDF: {file_path.name}")
    
    full_text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            if not pdf.pages:
                return None, "Empty PDF"
            
            # Extract text from all pages
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    full_text += page_text + "\n"
            
            # Check if we got enough text (not a scanned PDF)
            if len(full_text.strip()) < 50:
                print("Scanned PDF detected - using OCR")
                all_text = []
                for page_num, page in enumerate(pdf.pages):
                    try:
                        img = page.to_image(resolution=150).original
                        temp_path = UPLOAD_DIR / f"{file_path.stem}_page{page_num+1}.png"
                        img.save(temp_path, format="PNG")
                        _, page_text = process_text_image(temp_path, get_analyzer())
                        if page_text and page_text != "No text detected in image":
                            all_text.append(page_text)
                        # Clean up temp file
                        if temp_path.exists():
                            temp_path.unlink()
                    except Exception as e:
                        print(f"Error processing page {page_num+1}: {e}")
                
                combined_text = "\n".join(all_text)
                if not combined_text.strip():
                    return None, "No text could be extracted from PDF"
                
                full_text = combined_text
            
            # Process text for PII (whether from direct extraction or OCR)
            analyzer = get_analyzer()
            anonymizer = get_anonymizer()
            
            # Analyze for PII
            analyzer_results = analyzer.analyze(
                text=full_text, 
                language='en',
                entities=[
                    "PERSON", "EMAIL_ADDRESS", "PHONE_NUMBER", "LOCATION",
                    "CREDIT_CARD", "US_SSN", "US_DRIVER_LICENSE", "US_PASSPORT",
                    "DATE_TIME", "NRP", "MEDICAL_LICENSE", "URL",
                    "AWS_IDENTIFIER", "NETWORK_INFO", "FILE_PATH", "TIMESTAMP",
                    "IDENTIFIER_NAME", "DURATION"
                ]
            )
            
            # Anonymize the text
            anonymized_result = anonymizer.anonymize(text=full_text, analyzer_results=analyzer_results)
            
            # Save redacted version
            redacted_filename = f"{file_path.stem}_redacted.txt"
            redacted_path = REDACTED_DIR / redacted_filename
            with open(redacted_path, 'w', encoding='utf-8') as f:
                f.write("PDF Content - Redacted Version\n")
                f.write("=" * 40 + "\n\n")
                f.write(anonymized_result.text)
            
            print(f"PDF processing completed. Text length: {len(full_text)}")
            # Return original text for analysis (not anonymized version)
            return redacted_path, full_text
                
    except Exception as e:
        print(f"Error processing PDF: {e}")
        return None, f"Error processing PDF: {str(e)}"

def process_excel(file_path: Path):
    """Process Excel files (.xlsx, .xls)"""
    print(f"Processing Excel file: {file_path.name}")
    
    try:
        # Load the workbook
        if file_path.suffix.lower() == '.xls':
            # For .xls files, we need to use pandas or xlrd
            import pandas as pd
            excel_data = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
            
            all_text = []
            for sheet_name, df in excel_data.items():
                sheet_text = f"Sheet: {sheet_name}\n"
                # Convert DataFrame to text
                for index, row in df.iterrows():
                    row_text = []
                    for value in row:
                        if pd.notna(value):
                            row_text.append(str(value))
                    if row_text:
                        sheet_text += " | ".join(row_text) + "\n"
                all_text.append(sheet_text)
        else:
            # For .xlsx files, use openpyxl
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            
            all_text = []
            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"Sheet: {sheet_name}\n"
                
                # Extract text from cells
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        sheet_text += " | ".join(row_text) + "\n"
                
                all_text.append(sheet_text)
        
        # Combine all sheets
        full_text = "\n\n".join(all_text)
        
        if not full_text.strip():
            return None, "No data found in Excel file"
        
        # Analyze for PII
        analyzer = get_analyzer()
        anonymizer = get_anonymizer()
        
        analyzer_results = analyzer.analyze(
            text=full_text,
            language='en',
            entities=[
                "PERSON", "EMAIL_ADDRESS", "PHONE_NUMBER", "LOCATION",
                "CREDIT_CARD", "US_SSN", "US_DRIVER_LICENSE", "US_PASSPORT",
                "DATE_TIME", "NRP", "MEDICAL_LICENSE", "URL",
                "AWS_IDENTIFIER", "NETWORK_INFO", "FILE_PATH", "TIMESTAMP",
                "IDENTIFIER_NAME", "DURATION"
            ]
        )
        
        # Anonymize the text
        anonymized_result = anonymizer.anonymize(text=full_text, analyzer_results=analyzer_results)
        
        # Save redacted version
        redacted_filename = f"{file_path.stem}_redacted.txt"
        redacted_path = REDACTED_DIR / redacted_filename
        with open(redacted_path, 'w', encoding='utf-8') as f:
            f.write("Excel File Analysis - Redacted Content\n")
            f.write("=" * 50 + "\n\n")
            f.write(anonymized_result.text)
        
        print(f"Excel processing completed. Text length: {len(full_text)}")
        return redacted_path, full_text
        
    except Exception as e:
        print(f"Error processing Excel file: {e}")
        return None, f"Error processing Excel file: {str(e)}"

# Update the process_powerpoint function
def process_powerpoint(file_path: Path):
    """Process PowerPoint files (.ppt, .pptx) with slide image extraction"""
    print(f"Processing PowerPoint file: {file_path.name}")
    
    try:
        # Load the presentation
        prs = Presentation(file_path)
        
        all_text = []
        slide_count = 0
        slide_images = []
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_count += 1
            slide_text = f"Slide {slide_num}:\n"
            slide_content = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                
                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_content.append(" | ".join(row_text))
            
            # Add slide content
            if slide_content:
                slide_text += "\n".join(slide_content) + "\n"
                
                # Analyze slide content with Gemini for better description
                slide_description = analyze_slide_content_with_gemini(slide_content, slide_num)
                slide_text += f"Description: {slide_description}\n"
            else:
                slide_text += "[Slide contains visual elements or no readable text]\n"
                # Try to describe visual slide with Gemini
                slide_description = f"Slide {slide_num} contains primarily visual content with minimal text elements."
                slide_text += f"Description: {slide_description}\n"
            
            all_text.append(slide_text)
        
        # Try to extract slide images (this requires additional processing)
        slide_images_info = extract_slide_images(file_path, prs)
        
        # Combine all slides
        full_text = "\n\n".join(all_text)
        
        if not full_text.strip() or len(full_text.strip()) < 50:
            # If very little text, still create a basic analysis
            full_text = f"PowerPoint presentation with {slide_count} slides. Contains primarily visual content with minimal text elements."
        
        # Add presentation metadata
        presentation_info = f"PowerPoint Presentation Analysis:\nTotal Slides: {slide_count}\n\n{full_text}"
        
        if slide_images_info:
            presentation_info += f"\n\nSlide Images: {len(slide_images_info)} slides processed for visual content"
        
        # Analyze for PII
        analyzer = get_analyzer()
        anonymizer = get_anonymizer()
        
        analyzer_results = analyzer.analyze(
            text=full_text,
            language='en',
            entities=[
                "PERSON", "EMAIL_ADDRESS", "PHONE_NUMBER", "LOCATION",
                "CREDIT_CARD", "US_SSN", "US_DRIVER_LICENSE", "US_PASSPORT",
                "DATE_TIME", "NRP", "MEDICAL_LICENSE", "URL",
                "AWS_IDENTIFIER", "NETWORK_INFO", "FILE_PATH", "TIMESTAMP",
                "IDENTIFIER_NAME", "DURATION"
            ]
        )
        
        # Anonymize the text
        anonymized_result = anonymizer.anonymize(text=presentation_info, analyzer_results=analyzer_results)
        
        # Save redacted version with slide information
        redacted_filename = f"{file_path.stem}_redacted.txt"
        redacted_path = REDACTED_DIR / redacted_filename
        with open(redacted_path, 'w', encoding='utf-8') as f:
            f.write("PowerPoint Presentation - Redacted Content\n")
            f.write("=" * 55 + "\n\n")
            f.write(f"Original file: {file_path.name}\n")
            f.write(f"Total slides: {slide_count}\n")
            f.write(f"Processing date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            f.write("REDACTED CONTENT:\n")
            f.write("-" * 30 + "\n")
            f.write(anonymized_result.text)
            
            if slide_images_info:
                f.write(f"\n\nSlide Images Processed: {len(slide_images_info)}\n")
                for img_info in slide_images_info:
                    f.write(f"- {img_info}\n")
        
        print(f"PowerPoint processing completed. Slides: {slide_count}, Text length: {len(presentation_info)}")
        return redacted_path, presentation_info
        
    except Exception as e:
        print(f"Error processing PowerPoint file: {e}")
        return None, f"Error processing PowerPoint file: {str(e)}"

def analyze_slide_content_with_gemini(slide_content: list, slide_num: int):
    """Use Gemini to analyze and describe slide content"""
    try:
        gemini_model = get_gemini_model()
        if gemini_model is None or not slide_content:
            return f"Slide {slide_num} contains text and visual elements."
        
        content_text = "\n".join(slide_content)
        
        prompt = f"""Analyze this PowerPoint slide content and provide a brief description:

Slide {slide_num} Content:
{content_text}

Provide a 1-2 sentence description of what this slide is about, its main topic or purpose. Be concise and specific."""
        
        response = gemini_model.generate_content(prompt)
        description = response.text.strip()
        
        if not description or len(description) < 10:
            return f"Slide {slide_num} contains information about the presented topic."
        
        return description
        
    except Exception as e:
        print(f"Error analyzing slide content with Gemini: {e}")
        return f"Slide {slide_num} contains structured content and information."

def extract_slide_images(file_path: Path, presentation):
    """Extract images from PowerPoint slides"""
    try:
        slide_images = []
        
        # Note: Direct slide-to-image conversion requires additional libraries
        # For now, we'll create a manifest of slide content
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_info = {
                'slide_number': slide_num,
                'text_elements': 0,
                'image_elements': 0,
                'table_elements': 0,
                'chart_elements': 0
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_info['text_elements'] += 1
                elif shape.shape_type == 13:  # Picture type
                    slide_info['image_elements'] += 1
                elif hasattr(shape, 'has_table') and shape.has_table:
                    slide_info['table_elements'] += 1
                elif hasattr(shape, 'has_chart') and shape.has_chart:
                    slide_info['chart_elements'] += 1
            
            slide_description = f"Slide {slide_num}: {slide_info['text_elements']} text elements, {slide_info['image_elements']} images, {slide_info['table_elements']} tables, {slide_info['chart_elements']} charts"
            slide_images.append(slide_description)
        
        return slide_images
        
    except Exception as e:
        print(f"Error extracting slide images: {e}")
        return []

# Add this enhanced function for slide image extraction (optional)
def extract_slide_images_advanced(file_path: Path, presentation):
    """Extract actual slide images (advanced feature)"""
    try:
        import io
        from PIL import Image as PILImage
        
        slide_images = []
        
        # This is a simplified approach - full implementation would require 
        # additional libraries like python-pptx with image export capabilities
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            # Create slide summary instead of actual image for now
            slide_summary = analyze_slide_layout(slide, slide_num)
            
            # Save slide summary as a text "image" description
            slide_desc_file = REDACTED_DIR / f"{file_path.stem}_slide_{slide_num}_description.txt"
            with open(slide_desc_file, 'w', encoding='utf-8') as f:
                f.write(f"Slide {slide_num} Layout Analysis\n")
                f.write("=" * 30 + "\n\n")
                f.write(slide_summary)
            
            slide_images.append({
                'slide_number': slide_num,
                'description_file': slide_desc_file,
                'summary': slide_summary
            })
        
        return slide_images
        
    except Exception as e:
        print(f"Error in advanced slide extraction: {e}")
        return []

def analyze_slide_layout(slide, slide_num):
    """Analyze the layout and structure of a slide"""
    try:
        layout_info = {
            'title': '',
            'text_content': [],
            'visual_elements': [],
            'layout_type': 'Unknown'
        }
        
        # Extract title and content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.name.startswith('Title') or 'title' in shape.name.lower():
                    layout_info['title'] = shape.text.strip()
                else:
                    layout_info['text_content'].append(shape.text.strip())
            
            # Identify visual elements
            if shape.shape_type == 13:  # Picture
                layout_info['visual_elements'].append('Image/Picture')
            elif hasattr(shape, 'has_table') and shape.has_table:
                layout_info['visual_elements'].append('Table')
            elif hasattr(shape, 'has_chart') and shape.has_chart:
                layout_info['visual_elements'].append('Chart/Graph')
        
        # Determine layout type
        if layout_info['title'] and layout_info['text_content']:
            layout_info['layout_type'] = 'Title and Content'
        elif layout_info['title'] and not layout_info['text_content']:
            layout_info['layout_type'] = 'Title Only'
        elif layout_info['visual_elements']:
            layout_info['layout_type'] = 'Visual Content'
        else:
            layout_info['layout_type'] = 'Custom Layout'
        
        # Create summary
        summary = f"Layout Type: {layout_info['layout_type']}\n"
        if layout_info['title']:
            summary += f"Title: {layout_info['title']}\n"
        
        if layout_info['text_content']:
            summary += f"Text Elements: {len(layout_info['text_content'])}\n"
            for i, text in enumerate(layout_info['text_content'][:3], 1):
                summary += f"  {i}. {text[:100]}{'...' if len(text) > 100 else ''}\n"
        
        if layout_info['visual_elements']:
            summary += f"Visual Elements: {', '.join(layout_info['visual_elements'])}\n"
        
        return summary
        
    except Exception as e:
        return f"Slide {slide_num}: Layout analysis failed - {str(e)}"

# Update the format_technical_analysis function to handle PowerPoint content better
def format_technical_analysis(analysis_text: str, file_path: Path = None):
    """Format technical analysis into table structure using Gemini"""
    print("Formatting analysis with Gemini...")
    
    gemini_model = get_gemini_model()
    file_info = extract_file_info_from_path(file_path)
    
    if (gemini_model is None):
        return format_fallback_table(analysis_text, file_info)
    
    # Truncate very long text for API limits
    analysis_preview = analysis_text[:3000] if len(analysis_text) > 3000 else analysis_text
    
    # Detect file type for better prompting
    file_type = file_info['type'].lower()
    
    if file_type in ['.ppt', '.pptx']:
        content_type = "PowerPoint presentation"
        focus_areas = "slide content, presentation topics, data types, and key themes"
    elif file_type == '.pdf':
        content_type = "PDF document"
        focus_areas = "document structure, content themes, data categories, and information types"
    elif file_type in ['.xlsx', '.xls']:
        content_type = "Excel spreadsheet"
        focus_areas = "data structure, column types, numerical data, and organizational patterns"
    else:
        content_type = "document or image"
        focus_areas = "content type, visible elements, and key characteristics"
    
    prompt = f"""Analyze this {content_type} content and create a structured table. 

Content to analyze:
{analysis_preview}

File: {file_info['name']}{file_info['type']}

Create a markdown table with these exact columns: File Name | File Type | File Description | Key Findings

Instructions:
1. File Description: Write 2-3 sentences describing what this {content_type} contains, its purpose, or main topic
2. Key Findings: Create 4-5 bullet points focusing on {focus_areas}
3. Use • symbol with <br> between findings
4. Be specific and informative based on the actual content
5. For presentations, mention slide count and main topics
6. For documents, focus on structure and content themes
7. For data files, highlight data types and organization

Format exactly like this:
| File Name | File Type | File Description | Key Findings |
|-----------|-----------|------------------|--------------|
| {file_info['name']} | {file_info['type']} | [Your description here] | • [Finding 1]<br>• [Finding 2]<br>• [Finding 3]<br>• [Finding 4] |

Only output the table, nothing else."""

    try:
        response = gemini_model.generate_content(prompt)
        result = response.text.strip()
        
        # Clean up the response to ensure it's a proper table
        if not result.startswith('|'):
            # Try to extract table from response
            lines = result.split('\n')
            table_lines = [line for line in lines if line.strip().startswith('|')]
            if table_lines:
                result = '\n'.join(table_lines)
        
        print("Gemini formatting completed successfully")
        return result
    except Exception as e:
        print(f"Gemini formatting failed: {e}")
        return format_fallback_table(analysis_text, file_info)

def format_fallback_table(analysis_text: str, file_info: dict):
    """Create basic table when Gemini fails"""
    print("Using fallback table formatting...")
    
    # Better text analysis for fallback
    text = analysis_text.strip()
    file_type = file_info['type'].lower()
    
    # Determine description based on content and file type
    if "Visual Content Analysis:" in text:
        # This is an image description
        description_start = text.find("Visual Content Analysis:") + len("Visual Content Analysis:")
        description_part = text[description_start:].split(".")[0].strip()
        description = f"Image file containing visual content. {description_part}."
        
        findings = [
            "Visual content detected without readable text",
            "No personally identifiable information found", 
            "Image suitable for display or documentation",
            "Content appears to be non-sensitive visual material"
        ]
    elif "PowerPoint Presentation Analysis:" in text:
        # This is a PowerPoint presentation
        lines = text.split('\n')
        slide_info = next((line for line in lines if "Total Slides:" in line), "")
        slide_count = slide_info.split(':')[-1].strip() if slide_info else "multiple"
        
        description = f"PowerPoint presentation containing {slide_count} slides with structured content and information."
        
        findings = [
            f"Presentation contains {slide_count} slides",
            "Structured content with text and visual elements",
            "Professional presentation format detected",
            "Content suitable for analysis and review"
        ]
        
        # Try to extract more specific findings from content
        if "Slide" in text:
            slide_count_actual = text.count("Slide ")
            findings[0] = f"Presentation contains {slide_count_actual} slides"
            
        if any(keyword in text.lower() for keyword in ['data', 'analysis', 'report']):
            findings.append("Contains data or analytical content")
        if any(keyword in text.lower() for keyword in ['chart', 'graph', 'table']):
            findings.append("Includes charts, graphs, or tabular data")
            
    elif file_type in ['.pdf']:
        description = "PDF document containing structured text and information for analysis."
        findings = [
            "Multi-format document with text content",
            "Structured information and data elements",
            "Content extracted successfully for analysis",
            "Document suitable for review and processing"
        ]
    elif file_type in ['.xlsx', '.xls']:
        description = "Excel spreadsheet containing structured data across multiple worksheets."
        findings = [
            "Spreadsheet with organized data structure",
            "Multiple data columns and rows detected",
            "Numerical and text data elements present",
            "Data suitable for analysis and processing"
        ]
    else:
        # Generic text content
        sentences = [s.strip() for s in text.replace('\n', ' ').split('.') if len(s.strip()) > 10]
        
        if sentences:
            description = f"{sentences[0]}. This file contains structured data and textual information."
            
            findings = []
            if len(sentences) > 1:
                findings.extend(sentences[1:4])
            
            # Add generic findings if needed
            while len(findings) < 4:
                generic_findings = [
                    "Contains structured textual data",
                    "Multiple data elements present",
                    "Information suitable for analysis",
                    "Document contains readable content"
                ]
                for gf in generic_findings:
                    if gf not in findings and len(findings) < 4:
                        findings.append(gf)
        else:
            description = "Document contains data and information for processing."
            findings = [
                "Structured content detected",
                "Data elements present",
                "Information processed successfully", 
                "Content suitable for analysis"
            ]
    
    # Ensure we have exactly 4 findings
    findings = findings[:4]
    while len(findings) < 4:
        findings.append("Additional content elements present")
    
    findings_text = "<br>".join([f"• {f}" for f in findings])
    
    table = f"""| File Name | File Type | File Description | Key Findings |
|-----------|-----------|------------------|--------------|
| {file_info['name']} | {file_info['type']} | {description} | {findings_text} |"""
    
    return table

def format_for_html_display(markdown_table: str):
    """Convert markdown table to styled HTML"""
    if not markdown_table or not markdown_table.strip():
        return "<p>No analysis available</p>"
    
    print(f"Converting markdown table to HTML...")
    
    html = """
<table class="analysis-table">
    <thead>
        <tr>
            <th style="width: 12%">File Name</th>
            <th style="width: 8%">File Type</th>
            <th style="width: 35%">File Description</th>
            <th style="width: 45%">Key Findings</th>
        </tr>
    </thead>
    <tbody>
"""
    
    # Parse markdown table
    lines = [l.strip() for l in markdown_table.split('\n') if l.strip()]
    
    for line in lines:
        if not line.startswith('|'):
            continue
        if '---' in line:  # Skip separator row
            continue
        if 'File Name' in line:  # Skip header row
            continue
        
        # Parse table cells
        cells = [cell.strip() for cell in line.split('|')]
        cells = [c for c in cells if c]  # Remove empty cells
        
        if len(cells) >= 4:
            file_name = cells[0]
            file_type = cells[1]
            description = cells[2]
            findings = cells[3]
            
            # Convert <br> separated bullets to HTML list
            findings_items = [f.strip().replace('•', '').strip() 
                            for f in findings.split('<br>') if f.strip()]
            findings_html = '<ul class="key-findings">' + \
                          ''.join([f'<li>{item}</li>' for item in findings_items]) + \
                          '</ul>'
            
            html += f"""
        <tr>
            <td><strong>{file_name}</strong></td>
            <td>{file_type}</td>
            <td>{description}</td>
            <td>{findings_html}</td>
        </tr>
"""
    
    html += """
    </tbody>
</table>
"""
    print(f"Generated HTML successfully")
    return html

def format_generic_output(analysis_text: str, file_path: Path = None):
    """Format generic content into table using Gemini"""
    return format_technical_analysis(analysis_text, file_path)

def format_final_output(analysis_text: str, file_path: Path = None, is_technical: bool = False):
    """Main formatting wrapper"""
    return format_technical_analysis(analysis_text, file_path)

def redact_text_in_image(file_path: Path, analyzer_engine: AnalyzerEngine = None):
    """Alias for backward compatibility"""
    if analyzer_engine is None:
        analyzer_engine = get_analyzer()
    return process_text_image(file_path, analyzer_engine)

# --- Create analyzer instance ---
analyzer = get_analyzer()

# --- Exports ---
__all__ = [
    'process_image',
    'process_pdf',
    'process_excel',
    'process_powerpoint',  # Add this
    'process_text_image',
    'redact_text_in_image',
    'format_technical_analysis',
    'format_generic_output',
    'format_final_output',
    'format_for_html_display',
    'extract_file_info_from_path',
    'analyzer',
    'UPLOAD_DIR',
    'REDACTED_DIR',
    'MODEL_DIR'
]